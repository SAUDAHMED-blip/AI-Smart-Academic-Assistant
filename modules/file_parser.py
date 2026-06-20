"""
file_parser.py
Universal document text extractor for HSA AI Academy.
Supports: PDF, DOCX, PPTX, XLSX, CSV, TXT, MD, JSON
"""

import os
import json
import csv
import io

# Heavy optional imports are done lazily inside each helper function
# to optimize startup performance and keep the project lightweight.


# Map of supported extensions → human labels
SUPPORTED_EXTENSIONS = {
    ".pdf":  "PDF Document",
    ".docx": "Word Document",
    ".doc":  "Word Document (Legacy)",
    ".pptx": "PowerPoint Presentation",
    ".ppt":  "PowerPoint (Legacy)",
    ".xlsx": "Excel Spreadsheet",
    ".xls":  "Excel (Legacy)",
    ".csv":  "CSV Spreadsheet",
    ".txt":  "Plain Text",
    ".md":   "Markdown",
    ".json": "JSON Data",
}


def get_supported_extensions() -> list:
    return list(SUPPORTED_EXTENSIONS.keys())


def is_supported(filename: str) -> bool:
    ext = os.path.splitext(filename.lower())[1]
    return ext in SUPPORTED_EXTENSIONS


def extract_text(file_stream, filename: str) -> str:
    """
    Extract plain text from a file stream.

    Args:
        file_stream: A file-like object (BytesIO or Werkzeug FileStorage stream).
        filename:    Original filename (used to detect format).

    Returns:
        Extracted text as a single string.

    Raises:
        ValueError: If the file type is unsupported or a required library is missing.
    """
    ext = os.path.splitext(filename.lower())[1]

    if ext == ".pdf":
        return _extract_pdf(file_stream)
    elif ext in (".docx", ".doc"):
        return _extract_docx(file_stream)
    elif ext in (".pptx", ".ppt"):
        return _extract_pptx(file_stream)
    elif ext in (".xlsx", ".xls"):
        return _extract_xlsx(file_stream)
    elif ext == ".csv":
        return _extract_csv(file_stream)
    elif ext in (".txt", ".md"):
        return _extract_text(file_stream)
    elif ext == ".json":
        return _extract_json(file_stream)
    else:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Supported formats: {', '.join(SUPPORTED_EXTENSIONS.keys())}"
        )


# ── Private extraction helpers ───────────────────────────────────────────────

def _extract_pdf(stream) -> str:
    try:
        import PyPDF2
    except ImportError:
        raise ValueError("PDF support requires PyPDF2. Run: pip install PyPDF2")
    reader = PyPDF2.PdfReader(stream)
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    if not pages:
        raise ValueError("The PDF appears to be scanned or image-only (no extractable text).")
    return "\n\n".join(pages)


def _extract_docx(stream) -> str:
    try:
        from docx import Document as DocxDocument
    except ImportError:
        raise ValueError("DOCX support requires python-docx. Run: pip install python-docx")
    doc = DocxDocument(stream)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    # Also extract table cells
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    paragraphs.append(text)
    if not paragraphs:
        raise ValueError("The Word document appears to contain no extractable text.")
    return "\n\n".join(paragraphs)


def _extract_pptx(stream) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ValueError("PPTX support requires python-pptx. Run: pip install python-pptx")
    prs = Presentation(stream)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text.strip())
        if len(slide_lines) > 1:
            slides_text.append("\n".join(slide_lines))
    if not slides_text:
        raise ValueError("The PowerPoint file appears to contain no extractable text.")
    return "\n\n".join(slides_text)


def _extract_xlsx(stream) -> str:
    try:
        import openpyxl
    except ImportError:
        raise ValueError("XLSX support requires openpyxl. Run: pip install openpyxl")
    wb = openpyxl.load_workbook(stream, data_only=True)
    all_sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(str(c) if c is not None else "" for c in row)
            if row_text.strip():
                rows.append(row_text)
        if rows:
            all_sheets.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
    if not all_sheets:
        raise ValueError("The Excel file appears to be empty.")
    return "\n\n".join(all_sheets)


def _extract_csv(stream) -> str:
    # Handle both bytes and text streams
    try:
        content = stream.read()
        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(content))
        rows = ["\t".join(row) for row in reader if any(cell.strip() for cell in row)]
        if not rows:
            raise ValueError("The CSV file appears to be empty.")
        return "\n".join(rows)
    except Exception as e:
        raise ValueError(f"Failed to parse CSV: {str(e)}")


def _extract_text(stream) -> str:
    content = stream.read()
    if isinstance(content, bytes):
        content = content.decode("utf-8", errors="replace")
    return content.strip()


def _extract_json(stream) -> str:
    content = stream.read()
    if isinstance(content, bytes):
        content = content.decode("utf-8", errors="replace")
    try:
        data = json.loads(content)
        # Pretty-print for readability by the AI
        return json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        # Return raw if invalid JSON
        return content.strip()
